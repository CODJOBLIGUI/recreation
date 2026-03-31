import io
import os
import re
from pathlib import Path

from django.conf import settings

_EASYOCR_READERS = {}

_ACRONYM_RE = re.compile(r"\b[A-Z0-9]{2,}\b")
_DIGIT_WORDS_FR = {
    "0": "zéro",
    "1": "un",
    "2": "deux",
    "3": "trois",
    "4": "quatre",
    "5": "cinq",
    "6": "six",
    "7": "sept",
    "8": "huit",
    "9": "neuf",
}


def _parse_tts_acronyms(raw_text):
    replacements = {}
    spell_set = set()
    if not raw_text:
        return replacements, spell_set
    for line in raw_text.splitlines():
        line = line.strip()
        if not line or line.startswith("#"):
            continue
        if "=" in line:
            key, value = line.split("=", 1)
            key = key.strip().upper()
            value = value.strip()
            if value:
                replacements[key] = value
            else:
                spell_set.add(key)
        else:
            spell_set.add(line.strip().upper())
    return replacements, spell_set


def _spell_token(token):
    parts = []
    for ch in token:
        if ch.isdigit():
            parts.append(_DIGIT_WORDS_FR.get(ch, ch))
        else:
            parts.append(ch)
    return " ".join(parts)


def normalize_tts_text(text, appearance=None, use_original=False):
    """
    Normalise les sigles et MAJUSCULES avant TTS.
    - remplace via dictionnaire (SIGLE=prononciation)
    - épelle les sigles courts (2-4) si activé
    - convertit les longs sigles en forme "mot" (ex: UNESCO -> Unesco)
    """
    if not text:
        return text
    if use_original:
        return text
    if appearance and not getattr(appearance, "tts_use_normalization", True):
        return text

    replacements, spell_set = _parse_tts_acronyms(
        getattr(appearance, "tts_acronyms", "")
    )
    spell_unknown = True
    if appearance is not None:
        spell_unknown = getattr(appearance, "tts_spell_unknown", True)

    def _replace(match):
        token = match.group(0)
        if token in replacements:
            return replacements[token]
        if token in spell_set or (spell_unknown and 2 <= len(token) <= 4):
            return _spell_token(token)
        if len(token) > 4:
            return token[0] + token[1:].lower()
        return token

    return _ACRONYM_RE.sub(_replace, text)


def _map_lang_to_ocr(lang_code):
    if not lang_code:
        return None
    code = (lang_code or "").lower()
    if code.startswith("fr"):
        return "fr"
    if code.startswith("en"):
        return "en"
    if code.startswith("ar"):
        return "ar"
    return None


def _detect_lang(text):
    try:
        from langdetect import detect
    except Exception:
        return None
    try:
        detected = detect(text)
    except Exception:
        return None
    if detected in {"fr", "en", "ar"}:
        return detected
    return None


def detect_tts_language(text, selected=None):
    """
    Détecte la langue pour TTS (fr/en/ar) si l'utilisateur choisit "auto".
    Sinon, on respecte la langue sélectionnée.
    """
    if selected and selected != "auto":
        return selected
    detected = _detect_lang(text or "")
    return detected or (selected if selected else "fr")


def estimate_pages_from_text(text):
    if not text:
        return 0
    words = len([w for w in text.split() if w.strip()])
    if words == 0:
        return 0
    # Estimation rapide : 300 mots ~ 1 page
    return max(1, int((words + 299) / 300))


def count_pages_for_file(file_field, language_hint=None, force_ocr=False):
    if not file_field:
        return 0
    name = file_field.name or ""
    ext = Path(name).suffix.lower()
    local_path = _ensure_local_path(file_field)

    if ext in {".pdf"}:
        try:
            from PyPDF2 import PdfReader
        except Exception:
            return 0
        with open(local_path, "rb") as f:
            reader = PdfReader(f)
            try:
                return len(reader.pages)
            except Exception:
                return 0

    if ext in {".pptx"}:
        try:
            from pptx import Presentation
        except Exception:
            return 0
        prs = Presentation(local_path)
        return len(prs.slides)

    if ext in {".xlsx"}:
        try:
            import openpyxl
        except Exception:
            return 0
        wb = openpyxl.load_workbook(local_path, data_only=True)
        return max(1, len(wb.worksheets))

    if ext in {".jpg", ".jpeg", ".png"}:
        return 1

    # Pour les autres formats, estimation via texte extrait
    try:
        text = extract_text_from_file(file_field, language_hint=language_hint, force_ocr=force_ocr)
    except Exception:
        text = ""
    return estimate_pages_from_text(text)

def _ensure_local_path(file_field):
    if hasattr(file_field, "path") and os.path.exists(file_field.path):
        return file_field.path
    tmp_dir = Path(settings.BASE_DIR) / "tmp"
    tmp_dir.mkdir(parents=True, exist_ok=True)
    suffix = Path(file_field.name).suffix or ".bin"
    tmp_path = tmp_dir / f"upload_{os.path.basename(file_field.name)}"
    if hasattr(file_field, "read") and hasattr(file_field, "seek"):
        data = file_field.read()
        try:
            file_field.seek(0)
        except Exception:
            pass
        with open(tmp_path, "wb") as dst:
            dst.write(data)
        return str(tmp_path)
    with file_field.open("rb") as src, open(tmp_path, "wb") as dst:
        dst.write(src.read())
    return str(tmp_path)


def extract_text_from_file(file_field, language_hint=None, force_ocr=False):
    if not file_field:
        return ""

    name = file_field.name or ""
    ext = Path(name).suffix.lower()
    local_path = _ensure_local_path(file_field)

    if ext in {".txt"}:
        with open(local_path, "rb") as f:
            return f.read().decode(errors="ignore")

    if ext in {".docx"}:
        try:
            from docx import Document
        except Exception as exc:
            raise RuntimeError("python-docx n'est pas installé.") from exc
        doc = Document(local_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text)

    if ext in {".pdf"}:
        try:
            from PyPDF2 import PdfReader
        except Exception as exc:
            raise RuntimeError("PyPDF2 n'est pas installé.") from exc
        with open(local_path, "rb") as f:
            reader = PdfReader(f)
            pages = [p.extract_text() or "" for p in reader.pages]
        text = "\n".join(pages).strip()
        if text and not force_ocr:
            return text
        # OCR fallback for scanned PDFs (extract embedded images)
        try:
            import easyocr
            from PIL import Image
            import numpy as np
        except Exception as exc:
            raise RuntimeError("OCR PDF indisponible (EasyOCR/Pillow manquant).") from exc
        selected = _map_lang_to_ocr(language_hint)
        detected = _detect_lang(text) if text else None
        ocr_lang = selected or detected or "fr"
        langs = (ocr_lang,)
        if langs not in _EASYOCR_READERS:
            _EASYOCR_READERS[langs] = easyocr.Reader(list(langs), gpu=False)
        reader_ocr = _EASYOCR_READERS[langs]
        texts = []
        for page in reader.pages:
            images = getattr(page, "images", []) or []
            if images:
                for img in images:
                    try:
                        img_data = img.data
                        image = Image.open(io.BytesIO(img_data))
                        img_arr = np.array(image)
                        results = reader_ocr.readtext(img_arr, detail=0, paragraph=True)
                        texts.extend(results)
                    except Exception:
                        continue

        return "\n".join(texts)

    if ext in {".jpg", ".jpeg", ".png"}:
        try:
            import easyocr
        except Exception as exc:
            raise RuntimeError("easyocr n'est pas installé.") from exc
        selected = _map_lang_to_ocr(language_hint)
        ocr_lang = selected or "fr"
        langs = (ocr_lang,)
        if langs not in _EASYOCR_READERS:
            _EASYOCR_READERS[langs] = easyocr.Reader(list(langs), gpu=False)
        reader = _EASYOCR_READERS[langs]
        results = reader.readtext(local_path, detail=0, paragraph=True)
        return "\n".join(results)

    if ext in {".pptx"}:
        try:
            from pptx import Presentation
        except Exception as exc:
            raise RuntimeError("python-pptx n'est pas installé.") from exc
        prs = Presentation(local_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(t for t in texts if t)

    if ext in {".xlsx"}:
        try:
            import openpyxl
        except Exception as exc:
            raise RuntimeError("openpyxl n'est pas installé.") from exc
        wb = openpyxl.load_workbook(local_path, data_only=True)
        texts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        texts.append(str(cell))
        return "\n".join(texts)

    if ext in {".epub"}:
        try:
            from ebooklib import epub
            from bs4 import BeautifulSoup
        except Exception as exc:
            raise RuntimeError("EbookLib ou beautifulsoup4 n'est pas installé.") from exc
        book = epub.read_epub(local_path)
        texts = []
        for item in book.get_items():
            if item.get_type() == epub.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                texts.append(soup.get_text(" ", strip=True))
        return "\n".join(texts)

    raise RuntimeError("Type de fichier non pris en charge.")


